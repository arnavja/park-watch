from pptx import Presentation
from pptx.util import Inches, Pt
import docx

# Update Docx
try:
    doc = docx.Document("outputs/Demo_Video_Script.docx")
    doc.add_heading("The True Cost of Congestion (Added via AI)", level=2)
    p = doc.add_paragraph("When we talk about traffic, we often just think about being annoyed. But we wanted to quantify the actual economic damage. To calculate the true cost of this congestion, our model factors in three key variables: the average time commuters lose while idling, the cost of wasted fuel, and the baseline hourly wages of the workforce in this area. We found that the hidden economic drain is staggering—costing the city millions annually.")
    doc.save("outputs/Demo_Video_Script.docx")
    print("Successfully updated Demo_Video_Script.docx")
except Exception as e:
    print(f"Error updating docx: {e}")

# Update Pptx
try:
    prs = Presentation("outputs/Park_Watch_Pitch.pptx")
    # Title and Content slide layout
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The True Cost of Congestion"
    
    tf = content.text_frame
    tf.text = "Congestion Cost = Time Lost + Fuel Wasted + Emissions Impact"
    
    p = tf.add_paragraph()
    p.text = "Based on localized vehicle idle times."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Factors in current fuel prices and average commuter wages."
    p2.level = 1
    
    p3 = tf.add_paragraph()
    p3.text = "Quantifies the hidden economic drain of gridlock."
    p3.level = 1
    
    prs.save("outputs/Park_Watch_Pitch.pptx")
    print("Successfully added a slide to Park_Watch_Pitch.pptx")
except Exception as e:
    print(f"Error updating pptx: {e}")
