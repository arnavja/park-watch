"""Build the Park-Watch pitch deck for Gridlock Hackathon prototype review."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent.parent
OUT_PATH = ROOT / "outputs" / "Park_Watch_Pitch.pptx"
SCREENS = ROOT / "outputs" / "screens"

# ── Brand palette (Bengaluru Traffic Police navy + alert orange)
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ORANGE = RGBColor(0xFF, 0x6B, 0x2B)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_bar(slide, prs, title, subtitle=None):
    # Top navy bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.18), prs.slide_width - Inches(1), Inches(0.7)
    )
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True; p.font.size = Pt(26); p.font.color.rgb = WHITE

    if subtitle:
        p2 = tx.text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12); p2.font.color.rgb = LIGHT


def add_bullet_box(slide, left, top, width, height, items, size=16):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(8)


def add_metric_card(slide, left, top, width, height, value, label,
                     accent=ORANGE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = accent
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(28); p1.font.bold = True; p1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(11); p2.font.color.rgb = GREY


def add_footer(slide, prs, text="Park-Watch · Gridlock Hackathon · Team Byte_me_kaar"):
    fb = slide.shapes.add_textbox(
        0, prs.slide_height - Inches(0.35),
        prs.slide_width, Inches(0.3),
    )
    p = fb.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9); p.font.color.rgb = GREY


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 1 · TITLE
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

    tx = s.shapes.add_textbox(Inches(1), Inches(2.4),
                                prs.slide_width - Inches(2), Inches(3))
    p = tx.text_frame.paragraphs[0]
    p.text = "PARK-WATCH"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE

    p2 = tx.text_frame.add_paragraph()
    p2.text = "AI-driven parking enforcement intelligence for Bengaluru."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20); p2.font.color.rgb = ORANGE

    p3 = tx.text_frame.add_paragraph()
    p3.text = "Team Byte_me_kaar  ·  Gridlock Hackathon @ Flipkart HQ × BTP"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(14); p3.font.color.rgb = LIGHT

    # ── 2 · THE PROBLEM (BTP's own words)
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Operational problem statement",
                  "Theme 1 — Poor visibility on parking-induced congestion")
    add_bullet_box(s, Inches(0.6), Inches(1.4), Inches(12), Inches(5), [
        "Enforcement is patrol-based and reactive — officers go where they've always gone.",
        "No heatmap exists of parking violations vs. their congestion impact.",
        "Hard to prioritize which zones deserve enforcement effort.",
        "Result: high-impact illegal parking continues; patrol effort is wasted on low-impact zones.",
    ], size=18)
    add_footer(s, prs)

    # ── 3 · THE DATASET
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Dataset overview",
                  "BTP violation records · Nov 2023 – Apr 2024 · Bengaluru-wide")
    w = Inches(2.9); gap = Inches(0.2)
    left = Inches(0.5); top = Inches(1.5)
    add_metric_card(s, left, top, w, Inches(1.4), "298,450", "violations recorded")
    add_metric_card(s, left + w + gap, top, w, Inches(1.4), "152 days", "5-month window")
    add_metric_card(s, left + 2*(w+gap), top, w, Inches(1.4), "54", "police stations")
    add_metric_card(s, left + 3*(w+gap), top, w, Inches(1.4), "169", "BTP junctions")
    add_bullet_box(s, Inches(0.6), Inches(3.3), Inches(12), Inches(3.5), [
        "Every record: geo-tagged (lat/lon), timestamped, vehicle type, violation type, station, junction.",
        "~100% are parking-related: WRONG PARKING (27K), NO PARKING (23K), MAIN ROAD (4K), FOOTPATH (616)…",
        "Bounding box: lat 12.8 – 13.3, lon 77.4 – 77.8 — full GBA/BBMP area.",
        "Vehicle mix: Scooter (95K), Car (89K), Motorcycle (41K), Passenger Auto (38K).",
        "Scale context: 1.23 Cr registered vehicles · 14.4M population · ~4,792 traffic officers · ratio 1:2,303 vs BPR&D recommended 1:700.",
    ], size=13)
    add_footer(s, prs)

    # ── 4 · THE KILLER INSIGHT  (with screenshot of blind-spot chart)
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Enforcement coverage gap",
                  "Temporal analysis of 298K bookings reveals a 12-hour visibility window")
    add_bullet_box(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), [
        "~95% of 298K bookings happen before 3 PM, peaking 8 AM – noon.",
        "7–11 PM (peak commercial activity): only ~27, 42, 148, 725 violations booked across five months.",
        "This is not because illegal parking stops. It's because officers go home.",
        "Every evening, thousands of incidents are invisible to enforcement.",
        "Park-Watch closes this gap.",
    ], size=14)
    # Screenshot of the live dashboard's bar chart on the right
    img = SCREENS / "01_kpis_blindspot.png"
    if img.exists():
        s.shapes.add_picture(
            str(img), Inches(6.5), Inches(1.2),
            width=Inches(6.5), height=Inches(4.5),
        )
        cap = s.shapes.add_textbox(Inches(6.5), Inches(5.8),
                                     Inches(6.5), Inches(0.4))
        cp = cap.text_frame.paragraphs[0]
        cp.text = "Live Park-Watch dashboard — hourly violation distribution"
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(10); cp.font.italic = True; cp.font.color.rgb = GREY
    add_footer(s, prs)

    # ── 5 · THE SOLUTION ARCHITECTURE
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "System architecture",
                  "Four analytical modules integrated into a single operational dashboard")
    add_bullet_box(s, Inches(0.6), Inches(1.4), Inches(12), Inches(5), [
        "1.  HOTSPOT CLUSTERING — DBSCAN on 298K geo-points → 381 illegal-parking zones across BLR.",
        "2.  CONGESTION COST — OSMnx road graph + BPR delay model → ₹/month lost at each zone.",
        "3.  BLIND-SPOT FORECAST — XGBoost on 1.38M hourly cells → predicts where evening violations will happen.",
        "4.  PATROL OPTIMIZER — Weighted KMeans + nearest-neighbor TSP → tonight's optimal patrol routes.",
        "5.  DASHBOARD — Streamlit + Folium, with live patrol-count slider for BTP shift planning.",
        "6.  FEEDBACK LOOP — Daily incremental + weekly full retrain + monthly re-cluster. Calendar features capture annual patterns (Diwali, Ganesh Chaturthi, monsoon, school holidays).",
    ], size=16)
    add_footer(s, prs)

    # ── 6 · MODULE 1+2 RESULTS (hotspots + cost)
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Hotspot detection and congestion cost",
                  "Module 1 — DBSCAN clustering   ·   Module 2 — OSMnx + BPR delay model")
    add_metric_card(s, Inches(0.5), Inches(1.3), Inches(3), Inches(1.4),
                     "381", "hotspots discovered")
    add_metric_card(s, Inches(3.7), Inches(1.3), Inches(3), Inches(1.4),
                     "1.6 M", "vehicle-hours lost/month")
    add_metric_card(s, Inches(6.9), Inches(1.3), Inches(3), Inches(1.4),
                     "₹389 Cr", "lost per year")
    add_metric_card(s, Inches(10.1), Inches(1.3), Inches(2.9), Inches(1.4),
                     "₹16.4 Cr", "KR Market alone, monthly")
    add_bullet_box(s, Inches(0.5), Inches(3.0), Inches(5.5), Inches(4), [
        "Top 5: KR Market · Dispensary Rd (Shivajinagar) · Dr Rajkumar Rd · HAL Outer Ring · 10th Cross Malleshwaram.",
        "OSMnx pulls real lane count + speed limit per road segment.",
        "BPR (α=0.15, β=4) translates lane-blockage into delay.",
        "Value-of-time: ₹200/hr (GoK Economic Survey).",
        "₹389 Cr/yr ≈ 10% of BLR's ₹38,000 Cr congestion bill — from 100 parking zones alone.",
    ], size=12)
    img = SCREENS / "05_patrol_routes.png"
    if img.exists():
        s.shapes.add_picture(
            str(img), Inches(6.5), Inches(3.0),
            width=Inches(6.5), height=Inches(4),
        )
        cap = s.shapes.add_textbox(Inches(6.5), Inches(7.05),
                                     Inches(6.5), Inches(0.3))
        cp = cap.text_frame.paragraphs[0]
        cp.text = "Live dashboard — hotspot priority table"
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(10); cp.font.italic = True; cp.font.color.rgb = GREY
    add_footer(s, prs)

    # ── 7 · MODULE 3 — THE AI
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Violation forecasting model",
                  "Module 3 — XGBoost regression on 1.38M (hotspot × hour) observations")
    add_metric_card(s, Inches(0.5), Inches(1.3), Inches(3), Inches(1.4),
                     "1.38 M", "training observations")
    add_metric_card(s, Inches(3.7), Inches(1.3), Inches(3), Inches(1.4),
                     "0.43", "test R²  (chronological split)")
    add_metric_card(s, Inches(6.9), Inches(1.3), Inches(3), Inches(1.4),
                     "+0.07", "Train→Test gap (no overfit)")
    add_metric_card(s, Inches(10.1), Inches(1.3), Inches(2.9), Inches(1.4),
                     "743 / day", "unbooked violations forecast")
    add_bullet_box(s, Inches(0.5), Inches(3.0), Inches(5.5), Inches(4), [
        "Features (11): hour, dow, month, weekend, lag-1h/-24h/-7d, roll-7d, lat, lon, vehicle.",
        "Top importances: lag-1h (34%), roll-7d (16%), lag-24h (10%).",
        "70/15/15 chronological split, early stopping at iter 59 of 1500.",
        "Train R² 0.50 → Test R² 0.43 → gap +0.07 (≪ 0.15 threshold) → no overfit.",
        "Forecast: 24h per hotspot, blind-spot hours aggregated for patrol planning.",
        "Limitation: dataset captures booked violations, not all events. Forecast is a testable hypothesis a 1-week BTP pilot can validate.",
    ], size=12)
    img = SCREENS / "03_forecast.png"
    if img.exists():
        s.shapes.add_picture(
            str(img), Inches(6.5), Inches(3.0),
            width=Inches(6.5), height=Inches(4),
        )
        cap = s.shapes.add_textbox(Inches(6.5), Inches(7.05),
                                     Inches(6.5), Inches(0.3))
        cp = cap.text_frame.paragraphs[0]
        cp.text = "Live dashboard — 24h forecast for KR Market hotspot"
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(10); cp.font.italic = True; cp.font.color.rgb = GREY
    add_footer(s, prs)

    # ── 8 · MODULE 4 — PATROL PLAN
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Patrol route optimization",
                  "Module 4 — Weighted KMeans assignment + nearest-neighbor TSP sequencing")
    add_metric_card(s, Inches(0.5), Inches(1.3), Inches(3), Inches(1.4),
                     "5", "patrols deployed")
    add_metric_card(s, Inches(3.7), Inches(1.3), Inches(3), Inches(1.4),
                     "53", "optimized stops (5-hr shift)")
    add_metric_card(s, Inches(6.9), Inches(1.3), Inches(3), Inches(1.4),
                     "~96", "expected catches tonight")
    add_metric_card(s, Inches(10.1), Inches(1.3), Inches(2.9), Inches(1.4),
                     "~3.8×", "vs current BTP output (~25)")
    add_bullet_box(s, Inches(0.5), Inches(3.0), Inches(5.5), Inches(4), [
        "Patrols home-base: 5 real BTP stations — Upparpet, Shivajinagar, Malleshwaram, HAL Old Airport, Vijayanagara.",
        "Weighted KMeans → fair workload assignment by expected catches.",
        "Nearest-neighbor TSP within each patrol minimizes travel time.",
        "20 km/h urban speed + 15 min service per zone.",
        "Patrol count is a parameter — demo shows 5; in deployment, scales to BTP's actual nightly fleet.",
        "Output: per-officer schedule with ETA, stop sequence, expected catches.",
        "Same officer-hours, same budget. ~3.8× more enforcement output.",
    ], size=12)
    img = SCREENS / "04_hotspot_map.png"
    if img.exists():
        s.shapes.add_picture(
            str(img), Inches(6.5), Inches(3.0),
            width=Inches(6.5), height=Inches(4),
        )
        cap = s.shapes.add_textbox(Inches(6.5), Inches(7.05),
                                     Inches(6.5), Inches(0.3))
        cp = cap.text_frame.paragraphs[0]
        cp.text = "Live dashboard — color-coded patrol routes across BLR"
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(10); cp.font.italic = True; cp.font.color.rgb = GREY
    add_footer(s, prs)

    # ── 9 · WHY PARK-WATCH WINS
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Alignment with judging criteria")
    add_bullet_box(s, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5), [
        "APPLICATION OF TECHNOLOGY: 4 ML modules (DBSCAN, OSMnx+BPR, XGBoost, KMeans+TSP) integrated end-to-end into one decision system.",
        "BUSINESS VALUE: ₹389 Cr/year addressable inefficiency → ~3.8× enforcement output with zero extra cost → measurable BTP impact from day one.",
        "ORIGINALITY: Other teams will count violations. Only Park-Watch identifies BTP's structural enforcement blind spot and operationalizes the fix.",
        "PRESENTATION: Live Streamlit dashboard the BTP panel can navigate themselves — filter by station, vehicle, hour. No black box.",
        "DEPLOYABLE: Built entirely on the data BTP already collects. No new hardware. No new data pipeline. Just a dashboard.",
        "ROADMAP: v2 — adaptive per-officer routing. Mobile-app feed re-plans the route after every stop, learns each officer's pace, drops or adds zones based on remaining time.",
    ], size=15)
    add_footer(s, prs)

    # ── 10 · ARCHITECTURE DIAGRAM
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "End-to-end architecture",
                  "Data → analysis → decision  ·  6 modules, 1 dashboard, daily/weekly/monthly retrain")

    box_w = Inches(2.0); box_h = Inches(1.0); box_y = Inches(1.8)
    gap = Inches(0.15)
    arrows_y = Inches(2.3)
    pal = [
        ("BTP CSV\n298,450 rows", RGBColor(0x6B, 0x88, 0xB0)),
        ("Loader +\nParquet cache", RGBColor(0x6B, 0x88, 0xB0)),
        ("DBSCAN\n381 hotspots", ORANGE),
        ("OSMnx + BPR\n₹389 Cr/yr", ORANGE),
        ("XGBoost\nR²=0.43", ORANGE),
        ("KMeans + TSP\npatrol routes", ORANGE),
    ]
    x = Inches(0.4)
    for label, colour in pal:
        b = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, box_y, box_w, box_h
        )
        b.fill.solid(); b.fill.fore_color.rgb = colour
        b.line.fill.background()
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11); p.font.bold = True
        p.font.color.rgb = WHITE
        x += box_w + gap

    # downward arrow to dashboard
    dash = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.2),
        Inches(6.3), Inches(1.0),
    )
    dash.fill.solid(); dash.fill.fore_color.rgb = NAVY
    dash.line.fill.background()
    dtf = dash.text_frame
    dp = dtf.paragraphs[0]
    dp.text = "Streamlit + Folium dashboard"
    dp.alignment = PP_ALIGN.CENTER
    dp.font.size = Pt(16); dp.font.bold = True; dp.font.color.rgb = WHITE
    dp2 = dtf.add_paragraph()
    dp2.text = "Live patrol-count slider · per-officer shift sheet · TSP re-planner"
    dp2.alignment = PP_ALIGN.CENTER
    dp2.font.size = Pt(11); dp2.font.color.rgb = LIGHT

    # feedback loop arrow
    fb = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.6),
        Inches(12.3), Inches(0.7),
    )
    fb.fill.solid(); fb.fill.fore_color.rgb = LIGHT
    fb.line.color.rgb = ORANGE
    fb.line.width = Pt(1.5)
    fbtf = fb.text_frame
    fp = fbtf.paragraphs[0]
    fp.text = ("Production feedback loop: each shift's bookings retrain tomorrow's "
                "forecast · weekly full retrain · monthly DBSCAN re-cluster")
    fp.alignment = PP_ALIGN.CENTER
    fp.font.size = Pt(11); fp.font.color.rgb = NAVY; fp.font.italic = True

    add_footer(s, prs)

    # ── 11 · METHODOLOGY DEPTH
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Methodology — defensible knobs",
                  "Every number on the dashboard has explicit, auditable math behind it")

    # Left: BPR + cost model
    add_bullet_box(s, Inches(0.5), Inches(1.2), Inches(6.3), Inches(2.5), [
        "BPR delay model (Bureau of Public Roads, US-DOT standard):",
        "   t = t_free · (1 + α · (v/c)^β)    α=0.15  β=4",
        "Capacity loss per zone = viol/day × dwell_min ÷ (lanes × 1440)",
        "Value-of-time = ₹200/hr (GoK Economic Survey baseline)",
        "Cost cap: delay ratio clipped at 5× for conservative reporting",
    ], size=11)

    # Right: XGBoost hyperparameters
    add_bullet_box(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(2.5), [
        "XGBoost regressor — tuned for generalization:",
        "   n_estimators=1500   max_depth=6   lr=0.05",
        "   subsample=0.8   colsample_bytree=0.8   reg_lambda=1.0",
        "   min_child_weight=5   early_stopping=30 rounds",
        "Best iter=54  ·  test R²=0.43  ·  train→test gap +0.07",
    ], size=11)

    # Bottom: feature importance
    add_bullet_box(s, Inches(0.5), Inches(4.0), Inches(12.4), Inches(3), [
        "Feature importance (top 10) — recent activity + calendar dominate:",
        "   lag-1h 0.34  ·  roll-7d 0.16  ·  lag-7d 0.07  ·  hour 0.06  ·  lag-24h 0.05",
        "   top-vehicle-idx 0.03  ·  lat 0.03  ·  is-festival 0.03  ·  is-festival-window 0.03  ·  day-of-year 0.03",
        "Sample weights: 90-day exponential decay — recent data weighted higher, model adapts to evolving patrol patterns",
        "Calendar features: festivals (Diwali, Ganesh Chaturthi, Onam), monsoon (Jun-Sep), week-of-year, festival-window flags",
    ], size=12)
    add_footer(s, prs)

    # ── 12 · PER-OFFICER SHIFT SHEET
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Per-officer shift sheet — live tracking",
                  "v1 already delivers adaptive per-officer routing")

    add_bullet_box(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5), [
        "Each patrol gets a dropdown-selectable shift sheet with step-by-step orders.",
        "Step format: arrive time · depart time · travel min · expected catches · location.",
        "\"Mark stop done\" button advances state — remaining stops re-sequence live.",
        "Nearest-neighbour TSP replans from current position when an officer completes a stop.",
        "Time-budget guard: if replanned route overruns shift, system flags lowest-value stop to drop.",
        "Download button: shift sheet exports as plain TXT — printable, WhatsApp-able.",
    ], size=13)

    img = SCREENS / "05_patrol_routes.png"
    if img.exists():
        s.shapes.add_picture(
            str(img), Inches(6.5), Inches(1.2),
            width=Inches(6.5), height=Inches(5.0),
        )
        cap = s.shapes.add_textbox(Inches(6.5), Inches(6.3),
                                     Inches(6.5), Inches(0.3))
        cp = cap.text_frame.paragraphs[0]
        cp.text = "Live dashboard — per-officer shift sheet (Upparpet patrol)"
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(10); cp.font.italic = True; cp.font.color.rgb = GREY
    add_footer(s, prs)

    # ── 13 · PRODUCTION ROADMAP
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Production roadmap",
                  "v1 ships now — v2 adds personalized learning per officer")

    # 4-column timeline
    cols = [
        ("Day 0 — v1 deployment",
         "Streamlit dashboard live\nWeekly cron-scheduled retrain\nShift sheet exports to PDF/TXT/WhatsApp\nBTP officers patrol off the printed sheet"),
        ("Week 1–4 — pilot",
         "5 patrols × 5 hrs/night\nGround-truth catch rate measured\nForecast accuracy validated\nDataset doubles as evening data accumulates"),
        ("Month 2–3 — v1.5",
         "Mobile app for officers\nMark-stop-done in field\nLive route replan after each stop\nTime-budget warnings"),
        ("Month 4+ — v2",
         "Per-officer learning model\nIndividual pace + skill profile\nPersonalized routes per officer\nMulti-armed bandit zone exploration"),
    ]
    cw = Inches(3.0); ch = Inches(4.5); cy = Inches(1.3)
    cx = Inches(0.4)
    for title, body in cols:
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch
        )
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = ORANGE
        box.line.width = Pt(1.5)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.2); tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.2)
        p = tf.paragraphs[0]; p.text = title
        p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = ORANGE
        for line in body.split("\n"):
            p2 = tf.add_paragraph()
            p2.text = "•  " + line
            p2.font.size = Pt(11); p2.font.color.rgb = NAVY
            p2.space_after = Pt(4)
        cx += cw + Inches(0.15)
    add_footer(s, prs)

    # ── 14 · RISK / LIMITATIONS
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Limitations and risk register",
                  "Honest engineering — every assumption is auditable")

    add_bullet_box(s, Inches(0.5), Inches(1.3), Inches(12.5), Inches(5.5), [
        "DATASET CAPTURES BOOKINGS, NOT ALL VIOLATIONS — forecast extrapolates from booked patterns. A 1-week BTP pilot validates the forecast in the wild.",
        "PATROL FLEET + SHIFT WINDOW ARE PARAMETERS — demo runs 5 patrols in an evening window. BTP supplies actual roster — including the continuous towing patrols already running across 154 hotspots (DH, Apr 2026).",
        "JURISDICTIONAL SPLIT — BTP enforces, GBA/BBMP sets parking policy. Park-Watch supports both: enforcement priorities for BTP, hotspot analytics for GBA planners.",
        "BPR ASSUMPTIONS — dwell time (12 min), value-of-time (₹200/hr), lane capacity (1800 vph) are explicit knobs at the top of congestion_cost.py.",
        "OSM ROAD METADATA — lane counts and speed limits default to 2 lanes / 30 km/h when missing. Affects ~5% of edges in Bengaluru.",
        "EVENING DATA IS SPARSE — model trained on bookings from active enforcement hours. Predictions for 18:00–07:00 are extrapolations. Bootstraps once patrols arrive.",
        "MODEL DECAY RISK — without weekly retrain, accuracy drops as patrol patterns shift. Production cron-schedule prevents this.",
    ], size=13)
    add_footer(s, prs)

    # ── FINAL · CALL TO ACTION
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

    tx = s.shapes.add_textbox(Inches(1), Inches(2.5),
                                prs.slide_width - Inches(2), Inches(3))
    p = tx.text_frame.paragraphs[0]
    p.text = "Closing the enforcement coverage gap."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE

    p2 = tx.text_frame.add_paragraph()
    p2.text = "Same officers. Optimized routes. ~3.8× enforcement output."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(22); p2.font.color.rgb = ORANGE

    p3 = tx.text_frame.add_paragraph()
    p3.text = "park-watch.streamlit.app  ·  github.com/arnavja/park-watch  ·  arnavjain17032006@gmail.com"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(14); p3.font.color.rgb = LIGHT

    prs.save(OUT_PATH)
    print(f"Saved deck → {OUT_PATH}")


if __name__ == "__main__":
    build()
